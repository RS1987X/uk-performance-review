# agm/plotting.py

import matplotlib
matplotlib.use("Agg")     # <— ensure no GUI backend is used
import matplotlib.pyplot as plt

import pandas as pd
from pathlib import Path
from typing import Dict, Tuple

def save_performance_table(
    name: str,
    cum_pnl: pd.Series,
    drawdown: pd.Series,
    initial_capital: float,
    price_index_len: int,
    output_dir: Path
) -> Path:
    """
    Compute performance metrics, build a table figure, and save as PNG.
    Returns the path to the PNG.
    """
    from src.metrics import compute_performance_metrics

    perf_stats = compute_performance_metrics(
        cum_pnl=cum_pnl,
        drawdown=drawdown,
        initial_capital=initial_capital,
        trading_days=price_index_len
    )

    df = pd.DataFrame.from_dict(perf_stats, orient="index", columns=["Value"])
    df.index.name = "Metric"
    df = df.copy()
    df.loc["annualized_return", "Value"] = f"{df.loc['annualized_return','Value']*100:.2f}%"
    df.loc["sharpe_ratio", "Value"] = f"{df.loc['sharpe_ratio','Value']:.2f}"
    df.loc["gross_final_cum_pnl", "Value"] = f"{float(df.loc['gross_final_cum_pnl','Value']):.2f}"
    df.loc["max_drawdown", "Value"] = f"{float(df.loc['max_drawdown','Value']):.2f}"

    n_rows = df.shape[0] + 1
    row_height = 0.5
    fig_width = 6
    fig_height = n_rows * row_height

    fig, ax = plt.subplots(figsize=(fig_width, fig_height))
    ax.axis("off")
    tbl = ax.table(
        cellText=df.values,
        rowLabels=df.index,
        colLabels=df.columns,
        cellLoc="center",
        loc="center"
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(10)
    tbl.scale(1, 1.2)

    output_path = output_dir / f"Performance_stats_{name.replace(' ', '_')}.png"
    plt.savefig(output_path, bbox_inches="tight", dpi=150)
    plt.close(fig)
    return output_path

from pptx import Presentation
from pptx.util import Inches, Pt
import math
def save_df_as_ppt_table(df: pd.Series, output_path: Path, columns: int) -> None: 
    """
    Creates a PowerPoint slide with a compact table of PnL by ticker.
    Splits the DataFrame into the specified number of columns.
    Expects df to have columns ['ticker', 'final_pnl'] sorted descending.
    """
    # Ensure the Series is not empty
    if df.empty:
        raise ValueError("The input Series is empty.")

    # Convert the Series to lists for easier processing
    tickers = df.index.astype(str).tolist()
    pnl_values = df.values.tolist()

    # Determine layout
    total_items = len(df)
    rows_per_col = math.ceil(total_items / columns)
    table_cols = columns * 2  # ticker and PnL per column
    table_rows = rows_per_col + 1  # +1 for header row

    # Prepare Presentation and slide
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank slide

    # Table dimensions
    left = Inches(0.5)
    top = Inches(1.0)
    width = Inches(9.0)
    height = Inches(0.3 + 0.25 * rows_per_col)

    table = slide.shapes.add_table(table_rows, table_cols, left, top, width, height).table

    # Populate header
    for col in range(columns):
        table.cell(0, col * 2).text = "Ticker"
        table.cell(0, col * 2 + 1).text = "PnL"

    # Populate cells
    for idx in range(total_items):
        col = idx // rows_per_col
        row = idx % rows_per_col + 1  # +1 to skip header
        table.cell(row, col * 2).text = tickers[idx]
        table.cell(row, col * 2 + 1).text = f"{pnl_values[idx]:,.0f}"

    # Set font size for all cells
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)  # Set font size to 10 points

    # Save the presentation
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path / "pnl_by_ticker.pptx"))
    print(f"✅ Slide with PnL table saved to {output_path}")



def plot_portfolio_vs_omx(
    name: str,
    cum_pnl: pd.Series,
    drawdown: pd.Series,
    omx_cum_pnl: pd.Series,
    omx_drawdown: pd.Series,
    output_dir: Path
) -> Path:
    """
    Plot cumulative PnL + drawdown of a single portfolio vs. OMXSGI curves.
    Annotate peaks/troughs. Save to output_dir. Return path.
    """
    fig, ax = plt.subplots(figsize=(8, 4))
    ax.plot(cum_pnl, label="Cumulative PnL")
    ax.plot(drawdown, label="Drawdown")
    ax.plot(omx_cum_pnl, label="OMXSGI Equity", color="gray")
    ax.plot(omx_drawdown, label="OMXSGI Drawdown", color="lightgray", linestyle="--")

    high_date = cum_pnl.idxmax()
    high_value = cum_pnl.max()
    low_date = drawdown.idxmin()
    low_value = drawdown.min()

    ax.scatter([high_date], [high_value], color="green", zorder=5)
    ax.annotate(
        f"High: {high_value:.0f} SEK\n{high_date}",
        xy=(high_date, high_value),
        xytext=(10, 10),
        textcoords="offset points",
        arrowprops=dict(arrowstyle="->", color="green"),
        color="green"
    )
    ax.scatter([low_date], [low_value], color="red", zorder=5)
    ax.annotate(
        f"Low DD: {low_value:.0f} SEK\n{low_date}",
        xy=(low_date, low_value),
        xytext=(10, -30),
        textcoords="offset points",
        arrowprops=dict(arrowstyle="->", color="red"),
        color="red"
    )
    ax.set_title(name)
    ax.legend()
    plt.tight_layout()

    filename = f"{name.replace(' ', '_')}.png"
    output_path = output_dir / filename
    plt.savefig(output_path, dpi=150)
    plt.close(fig)
    return output_path


def plot_all_portfolios_vs_omx(
    results: Dict[str, Tuple[pd.Series, pd.Series]],
    omx_cum_pnl: pd.Series,
    output_dir: Path
) -> Path:
    """
    Plot each portfolio’s cum PnL plus an overlay of OMXSGI.
    """
    fig, ax = plt.subplots(figsize=(12, 6))
    for name, (cum_pnl, _) in results.items():
        ax.plot(cum_pnl, label=f"{name} Cum PnL", linewidth=1.5)

    ax.plot(omx_cum_pnl, label="OMXSGI Cum PnL", color="gray", linewidth=2)
    ax.set_title("All Portfolios: Cumulative PnL vs. OMXSGI", fontsize=14)
    ax.set_xlabel("Date", fontsize=12)
    ax.set_ylabel("Cumulative PnL (SEK)", fontsize=12)
    ax.legend(loc="upper left", bbox_to_anchor=(1.02, 1.00), fontsize=9)
    ax.grid(alpha=0.3)
    plt.tight_layout()

    output_path = output_dir / "all_portfolios_vs_omx_pnl.png"
    plt.savefig(output_path, bbox_inches="tight", dpi=150)
    plt.close(fig)
    print(f"Saved combined PnL plot → {output_path}")
    return output_path
